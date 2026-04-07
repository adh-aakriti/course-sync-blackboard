import pdfplumber
from pptx import Presentation


def extract_pdf(path):
    text = ""

    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text() or ""
                text += page_text + "\n"
    except Exception:
        print(f"Skipped invalid PDF: {path}")
        return ""

    return text.strip()


def extract_pptx(path):
    text = ""

    try:
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text += shape.text + "\n"
    except Exception:
        print(f"Skipped invalid PPTX: {path}")
        return ""

    return text.strip()


def extract_text(path):
    lower = path.lower()

    if lower.endswith(".pdf"):
        return extract_pdf(path)

    if lower.endswith(".pptx"):
        return extract_pptx(path)

    # old .ppt not supported directly here
    if lower.endswith(".ppt"):
        print(f"Skipped unsupported PPT format: {path}")
        return ""

    return ""